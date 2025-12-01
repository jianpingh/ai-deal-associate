from langchain_core.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from typing import List, Dict, Any

@tool
def create_presentation_slide(title: str, content: List[str], layout_index: int = 1) -> str:
    """
    Create a new PowerPoint presentation with a single slide (for demo purposes).
    In a real app, this might append to an existing deck or handle complex layouts.
    
    Args:
        title: The title of the slide.
        content: A list of bullet points or text content for the slide.
        layout_index: The index of the slide layout to use (default 1 is usually Title + Content).
    """
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        # Set content (assuming standard placeholder exists)
        # This is a simplified implementation
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1: # Usually the body text
                text_frame = shape.text_frame
                text_frame.clear()
                for point in content:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
        
        output_path = "generated_slide.pptx"
        prs.save(output_path)
        return f"Successfully created presentation at {output_path}"
    except Exception as e:
        return f"Error creating presentation: {str(e)}"
