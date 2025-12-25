from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_deal_summary_template():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "{{DEAL_NAME}}"
    subtitle.text = "Deal Summary / One-Pager\n{{DATE}}"

    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{SUMMARY_BULLETS}}"

    # Save
    base_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(base_dir, "data", "templates")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "deal_summary_template.pptx")
    prs.save(output_path)
    print(f"Deal Summary Template created at {output_path}")

def create_ic_deck_template():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "{{DEAL_NAME}}"
    subtitle.text = "Investment Committee Presentation\n{{DATE}}\n{{SCENARIO_LABEL}}"

    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{SUMMARY_BULLETS}}"

    # Slide 3: Market Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Market Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{MARKET_BULLETS}}"

    # Slide 4: Tenancy Schedule
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Tenancy Schedule"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{TENANCY_BULLETS}}"

    # Slide 5: Business Plan
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Business Plan"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{BUSINESS_PLAN_BULLETS}}"

    # Slide 6: Financial Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Financial Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Key Metrics:"
    p = tf.add_paragraph()
    p.text = "Entry Yield: {{ENTRY_YIELD}}"
    p = tf.add_paragraph()
    p.text = "IRR: {{IRR}}"
    p = tf.add_paragraph()
    p.text = "Equity Multiple: {{MOIC}}"
    p = tf.add_paragraph()
    p.text = "Exit Yield: {{EXIT_YIELD}}"

    # Slide 7: Sensitivities
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Sensitivities"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{SENSITIVITY_ANALYSIS}}"

    # Slide 8: Appendix
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Appendix"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "{{APPENDIX_BULLETS}}"

    # Save
    base_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(base_dir, "data", "templates")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "ic_deck_template.pptx")
    prs.save(output_path)
    print(f"IC Deck Template created at {output_path}")

if __name__ == "__main__":
    create_deal_summary_template()
    create_ic_deck_template()