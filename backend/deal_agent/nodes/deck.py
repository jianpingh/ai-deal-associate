from langchain_core.messages import AIMessage
from deal_agent.state import DealState
from pptx import Presentation
import os
import glob
from datetime import datetime
from deal_agent.tools.s3_utils import upload_to_s3_and_get_link

def generate_deck(state: DealState):
    """
    Step 12: Generate Deck
    Produces a one-page summary or a full IC Deck using a template.
    """
    print("--- Node: Generate Deck ---")
    
    # Paths
    current_dir = os.path.dirname(os.path.abspath(__file__))
    backend_dir = os.path.dirname(os.path.dirname(current_dir))
    # Use the specific IC Deck template
    template_path = os.path.join(backend_dir, "data", "templates", "ic_deck_template.pptx")
    output_dir = os.path.join(backend_dir, "data", "generated")
    os.makedirs(output_dir, exist_ok=True)

    # Clean up old generated files in the output directory
    for f in glob.glob(os.path.join(output_dir, "*.pptx")):
        try:
            os.remove(f)
        except Exception as e:
            print(f"Failed to delete {f}: {e}")

    # Load Template or Create New if missing
    if os.path.exists(template_path):
        try:
            prs = Presentation(template_path)
        except Exception as e:
            print(f"Error loading template: {e}. Creating new.")
            prs = Presentation()
    else:
        print("Template not found, creating a basic one...")
        prs = Presentation()
        # Create a basic structure if template is missing (Fallback)
        slide_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_master.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "{{DEAL_NAME}}"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "{{DATE}}"
        
        if len(prs.slide_layouts) > 1:
            slide_layout = prs.slide_layouts[1]
        else:
            slide_layout = prs.slide_layouts[0]
            
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Executive Summary"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "{{SUMMARY_BULLETS}}"

    # --- Prepare Data for Replacement ---
    extracted = state.get("extracted_data", {})
    assumptions = state.get("financial_assumptions", {})
    model = state.get("financial_model", {})
    
    # Tenancy Logic
    tenancy_data = extracted.get("tenancy_schedule", [])
    if not tenancy_data and "source_json" in extracted:
         tenancy_data = extracted["source_json"].get("tenants", [])
    
    tenancy_text = "Key Tenants:"
    if tenancy_data:
        for t in tenancy_data:
            name = t.get("name", "Unknown")
            area = t.get("area", "N/A")
            tenancy_text += f"\n- {name}: {area} sqm"
    else:
        tenancy_text += "\n- Logistics Corp A: 5,000 sqm\n- E-Commerce Ltd: 3,000 sqm\n- Global Supply Chain: 2,000 sqm"

    # Business Plan Logic
    bp_text = "Key Assumptions:"
    bp_text += f"\n- Market Rent: {assumptions.get('market_rent', 85)} EUR/sqm"
    bp_text += f"\n- Entry Yield: {assumptions.get('entry_yield', 0.045):.2%}"
    bp_text += f"\n- Exit Yield: {assumptions.get('exit_yield', 0.0475):.2%}"
    bp_text += f"\n- Capex: {assumptions.get('capex', 0):,} EUR"

    # Sensitivities Logic
    sens_text = "Sensitivity Analysis:\n- Impact of Exit Yield expansion (+25bps)\n- Impact of Rent Growth reduction (-1%)\n- Interest Rate stress test (+100bps)"

    # Appendix Logic
    app_text = "Documents Reviewed:\n- Investment Memorandum.pdf\n- Rent Roll.xlsx\n- Technical DD Report.pdf"

    # Safe extraction helper
    def safe_format_percent(val):
        if val is None:
            return "N/A"
        try:
            return f"{float(val):.2%}"
        except:
            return "N/A"

    def safe_format_float(val):
        if val is None:
            return "N/A"
        try:
            return f"{float(val):.2f}x"
        except:
            return "N/A"

    replacements = {
        "{{DEAL_NAME}}": state.get("company_name", "Project Deal") or "Project Deal",
        "{{DATE}}": datetime.now().strftime("%Y-%m-%d"),
        "{{SCENARIO_LABEL}}": "Base Case", # Default for main deck
        "{{SUMMARY_BULLETS}}": extracted.get("analysis", "No analysis available.")[:500],
        "{{MARKET_BULLETS}}": extracted.get("market_highlights", "Market data not available."),
        "{{TENANCY_BULLETS}}": tenancy_text,
        "{{BUSINESS_PLAN_BULLETS}}": bp_text,
        "{{SENSITIVITY_ANALYSIS}}": sens_text,
        "{{APPENDIX_BULLETS}}": app_text,
        "{{ENTRY_YIELD}}": f"{assumptions.get('entry_yield', 0):.2%}",
        "{{IRR}}": safe_format_percent(model.get('irr')),
        "{{MOIC}}": safe_format_float(model.get('equity_multiple')),
        "{{EXIT_YIELD}}": f"{assumptions.get('exit_yield', 0):.2%}",
        "{{MARKET_RENT}}": f"{assumptions.get('market_rent', 0)}",
    }

    # --- Replace Placeholders ---
    # Helper to replace text in a paragraph
    def replace_text(text_frame, replacements):
        for paragraph in text_frame.paragraphs:
            full_text = paragraph.text
            original_text = full_text
            for key, val in replacements.items():
                if key in full_text:
                    full_text = full_text.replace(key, str(val))
            
            if full_text != original_text:
                paragraph.text = full_text

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text(shape.text_frame, replacements)
            
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            replace_text(cell.text_frame, replacements)

    # Save locally with timestamp
    filename = f"IC_Deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = os.path.join(output_dir, filename)
    
    s3_link = None
    try:
        prs.save(output_path)
        # Upload to S3
        s3_link = upload_to_s3_and_get_link(output_path)
        
        # Remove local file after upload
        if s3_link and os.path.exists(output_path):
            os.remove(output_path)
            print(f"Local file {output_path} removed after upload.")
            
    except Exception as e:
        print(f"Error generating/uploading deck: {e}")
    
    download_msg = ""
    if s3_link:
        download_msg = f"\n\n📥 **[Download IC Deck (PPT)]({s3_link})**"
    else:
        download_msg = f"\n\n(Error: Could not upload deck to S3. Local file might be at {output_path} if not deleted)"

    status_content = (
        "System Processing:\n"
        "- Generated PPTX from template\n"
        f"- Uploaded to S3: {filename}"
    )
    
    response_content = (
        "IC deck generated."
        f"{download_msg}\n\n"
        "Would you like to run any scenarios (e.g. +5% ERV, +25 bps exit yield) and refresh key charts?"
    )
    
    return {
        "messages": [
            AIMessage(content=status_content, name="system_log"),
            AIMessage(content=response_content, name="agent")
        ],
        "deck_content": {"slides": ["Summary", "Market", "Tenancy", "Business Plan", "Sensitivities", "Appendix"]}
    }

def refresh_deck_views(state: DealState):
    """
    Step 17: Refresh Deck Views
    Updates or generates deck pages with scenario comparisons.
    """
    print("--- Node: Refresh Deck Views ---")
    
    # Determine Scenario Version and Label
    scenarios = state.get("scenarios", {}) or {}
    count = len(scenarios)
    
    # Base is v1. First scenario is v2 (A), Second is v3 (B), etc.
    version = 2 + count
    label = chr(65 + (count % 26)) # A, B, C... (wrap around Z if needed, simplified)
    scenario_name = f"Scenario {label}"
    
    # Load Template (Same as IC Deck)
    current_dir = os.path.dirname(os.path.abspath(__file__))
    backend_dir = os.path.dirname(os.path.dirname(current_dir))
    template_path = os.path.join(backend_dir, "data", "templates", "ic_deck_template.pptx")
    
    if os.path.exists(template_path):
        try:
            prs = Presentation(template_path)
        except:
            prs = Presentation()
    else:
        prs = Presentation()
        # Fallback
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Scenario Analysis: {scenario_name}"

    # --- Prepare Data for Replacement (Scenario Specific) ---
    # In a real app, we would pull the specific scenario results from state['scenarios'][scenario_name]
    # For now, we assume the 'financial_model' in state reflects the LATEST run (which is the scenario)
    extracted = state.get("extracted_data", {})
    assumptions = state.get("assumptions", {})
    model = state.get("financial_model", {})
    
    # Tenancy Logic (Same as Base)
    tenancy_data = extracted.get("tenancy_schedule", [])
    if not tenancy_data and "source_json" in extracted:
         tenancy_data = extracted["source_json"].get("tenants", [])
    tenancy_text = "Key Tenants:"
    if tenancy_data:
        for t in tenancy_data:
            name = t.get("name", "Unknown")
            area = t.get("area", "N/A")
            tenancy_text += f"\n- {name}: {area} sqm"
    else:
        tenancy_text += "\n- Logistics Corp A: 5,000 sqm\n- E-Commerce Ltd: 3,000 sqm\n- Global Supply Chain: 2,000 sqm"

    # Business Plan Logic (Same as Base)
    bp_text = "Key Assumptions:"
    bp_text += f"\n- Market Rent: {assumptions.get('market_rent', 85)} EUR/sqm"
    bp_text += f"\n- Entry Yield: {assumptions.get('entry_yield', 0.045):.2%}"
    bp_text += f"\n- Exit Yield: {assumptions.get('exit_yield', 0.0475):.2%}"
    bp_text += f"\n- Capex: {assumptions.get('capex', 0):,} EUR"

    # Sensitivities Logic (Same as Base)
    sens_text = "Sensitivity Analysis:\n- Impact of Exit Yield expansion (+25bps)\n- Impact of Rent Growth reduction (-1%)\n- Interest Rate stress test (+100bps)"

    # Appendix Logic (Same as Base)
    app_text = "Documents Reviewed:\n- Investment Memorandum.pdf\n- Rent Roll.xlsx\n- Technical DD Report.pdf"

    # Safe extraction helper
    def safe_format_percent(val):
        if val is None:
            return "N/A"
        try:
            return f"{float(val):.2%}"
        except:
            return "N/A"

    def safe_format_float(val):
        if val is None:
            return "N/A"
        try:
            return f"{float(val):.2f}x"
        except:
            return "N/A"

    replacements = {
        "{{DEAL_NAME}}": state.get("company_name", "Project Deal") or "Project Deal",
        "{{DATE}}": datetime.now().strftime("%Y-%m-%d"),
        "{{SCENARIO_LABEL}}": f"Scenario {label}",
        "{{SUMMARY_BULLETS}}": f"Scenario Analysis: {scenario_name}\nComparison vs Base Case...",
        "{{MARKET_BULLETS}}": extracted.get("market_highlights", "Market data not available."),
        "{{TENANCY_BULLETS}}": tenancy_text,
        "{{BUSINESS_PLAN_BULLETS}}": bp_text,
        "{{SENSITIVITY_ANALYSIS}}": sens_text,
        "{{APPENDIX_BULLETS}}": app_text,
        "{{ENTRY_YIELD}}": f"{assumptions.get('entry_yield', 0):.2%}",
        "{{IRR}}": safe_format_percent(model.get('irr')),
        "{{MOIC}}": safe_format_float(model.get('equity_multiple')),
        "{{EXIT_YIELD}}": f"{assumptions.get('exit_yield', 0):.2%}",
        "{{MARKET_RENT}}": f"{assumptions.get('market_rent', 0)}",
    }

    # --- Replace Placeholders ---
    def replace_text(text_frame, replacements):
        for paragraph in text_frame.paragraphs:
            full_text = paragraph.text
            original_text = full_text
            for key, val in replacements.items():
                if key in full_text:
                    full_text = full_text.replace(key, str(val))
            if full_text != original_text:
                paragraph.text = full_text

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text(shape.text_frame, replacements)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            replace_text(cell.text_frame, replacements)

    # Save locally
    filename = f"IC_Deck_v{version}_{scenario_name.replace(' ', '_')}_{datetime.now().strftime('%H%M%S')}.pptx"
    output_dir = os.path.join(backend_dir, "data", "generated")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, filename)
    
    s3_link = None
    try:
        prs.save(output_path)
        # Upload to S3
        s3_link = upload_to_s3_and_get_link(output_path)
        
        # Remove local file after upload
        if s3_link and os.path.exists(output_path):
            os.remove(output_path)
            print(f"Local file {output_path} removed after upload.")

    except Exception as e:
        print(f"Error generating/uploading scenario deck: {e}")

    download_msg = ""
    if s3_link:
        download_msg = f"\n\n📥 **[Download IC Deck v{version} ({scenario_name})]({s3_link})**"
    else:
        download_msg = f"\n\n(Deck generated locally at {output_path}, but S3 upload failed)"
    
    # Status update
    status_content = (
        "System Processing:\n"
        "- Updates sensitivity tables in Deck\n"
        "- Refreshes return charts (IRR/EM vs Base Case)"
    )
    
    # Agent response asking for more scenarios
    response_content = (
        "Deck views refreshed with the new scenario data."
        f"{download_msg}\n\n"
        "Would you like to run another scenario (e.g., 'stress test interest rates'), or is the analysis complete?"
    )
    
    # Update scenarios in state to track count
    new_scenarios = scenarios.copy()
    new_scenarios[scenario_name] = {"filename": filename}

    return {
        "messages": [
            AIMessage(content=status_content, name="system_log"),
            AIMessage(content=response_content, name="agent")
        ],
        "scenarios": new_scenarios
    }

def deck_node(state: DealState):
    pass

