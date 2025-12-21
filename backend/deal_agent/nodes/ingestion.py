import os
import json
from datetime import datetime
from langchain_core.messages import AIMessage, HumanMessage
from langchain_openai import ChatOpenAI
from deal_agent.state import DealState
from deal_agent.tools.pdf_parser import parse_pdf_document
from deal_agent.tools.vector_store import ingest_deal_assets
from pptx import Presentation
from deal_agent.tools.s3_utils import upload_to_s3_and_get_link

# --- Granular Nodes for Real-Time Logging ---

def start_ingestion(state: DealState):
    """
    Step 1: Start Ingestion
    Initializes the process and sends a log message.
    """
    print("--- Node: Start Ingestion ---")
    return {
        "messages": [
            AIMessage(content="Starting data ingestion process...", name="system_log")
        ]
    }

def load_json_data(state: DealState):
    """
    Step 2: Load JSON Data
    Loads the structured JSON file.
    """
    print("--- Node: Load JSON Data ---")
    
    base_dir = os.getcwd()
    if "backend" in base_dir:
        # Already in backend or subdirectory
        # If in backend/, data is in data/
        # If in backend/deal_agent/, data is in ../data/
        # Safest is to find 'backend' and construct path from there
        # But assuming standard execution from root or backend:
        if base_dir.endswith("backend"):
             data_root = os.path.join(base_dir, "data")
        else:
             # Fallback or deeper nesting, try to find relative to file
             current_file_dir = os.path.dirname(os.path.abspath(__file__))
             # .../backend/deal_agent/nodes -> .../backend
             backend_dir = os.path.dirname(os.path.dirname(current_file_dir))
             data_root = os.path.join(backend_dir, "data")
    else:
        # In root, data is in backend/data
        data_root = os.path.join(base_dir, "backend", "data")
        
    json_dir = os.path.join(data_root, "structured_json")
    structured_data = {}
    json_files = []
    
    if os.path.exists(json_dir):
        json_files = [f for f in os.listdir(json_dir) if f.endswith('.json')]
        if json_files:
            try:
                file_path = os.path.join(json_dir, json_files[0])
                with open(file_path, 'r', encoding='utf-8') as f:
                    structured_data = json.load(f)
                
                print(f"Vectorizing JSON data from {json_files[0]}...")
                deal_id = state.get("current_deal_id", "unknown_deal")
                
                # --- 1. Process Assets (Internal) -> Namespace: deal ---
                assets = structured_data.get("assets", [])
                asset_texts = []
                asset_metadatas = []
                
                for idx, asset in enumerate(assets):
                    # Extract nested fields
                    logistics = asset.get("logistics_asset", {})
                    leases = asset.get("leases", [])
                    tenant_names = [l.get("tenant", {}).get("name", "Unknown") for l in leases]
                    
                    # --- Currency & Unit Conversion Logic ---
                    # Original Data: GBP PSF (assumed for UK assets if value < 10)
                    # Target Data: EUR PSM
                    # Conversion Factors:
                    # 1. Unit: 1 sq m = 10.764 sq ft. So Price PSM = Price PSF * 10.764
                    # 2. Currency: 1 GBP = 1.2 EUR (Approx)
                    
                    raw_currency = asset.get('currency', 'GBP')
                    display_currency = "EUR"
                    conversion_note = ""
                    
                    # Calculate weighted average rent for display
                    total_rent = 0
                    total_area = 0
                    
                    for lease in leases:
                        raw_rent = lease.get('rent_psm_pa', 0) # In JSON this is actually PSF if it's ~5.5
                        area = lease.get('area_m2', 0)
                        
                        # Conversion Logic
                        if raw_currency == "GBP" and raw_rent < 20: # Heuristic check for PSF
                            # Convert GBP PSF -> EUR PSM
                            rent_eur_psm = raw_rent * 10.764 * 1.2
                            conversion_note = f"(Converted from £{raw_rent} PSF to €{rent_eur_psm:.2f} PSM)"
                        else:
                            rent_eur_psm = raw_rent # Assume already EUR PSM or close enough
                            
                        total_rent += rent_eur_psm * area
                        total_area += area
                        
                        # Update lease info for text blob (optional, but good for detail)
                        lease['display_rent_eur_psm'] = round(rent_eur_psm, 2)

                    avg_rent_eur_psm = round(total_rent / total_area, 2) if total_area > 0 else 0

                    text_blob = f"""
                    Asset Name: {asset.get('name', 'Unknown')}
                    Type: {asset.get('asset_type', 'Logistics')} ({asset.get('tenure', '')})
                    Location: {asset.get('address', '')}, {asset.get('city', '')}, {asset.get('country', '')}
                    Size: {logistics.get('area_m2', 0)} sqm
                    Specs: {logistics.get('eaves_height_m', 0)}m height, {logistics.get('dock_doors', 0)} docks
                    Tenants: {', '.join(tenant_names)}
                    Current Rent (Normalized): €{avg_rent_eur_psm} /sqm/year
                    Currency Basis: Converted to EUR PSM for comparison. {conversion_note}
                    """
                    asset_texts.append(text_blob.strip())
                    asset_metadatas.append({
                        "source": json_files[0],
                        "deal_id": deal_id,
                        "record_type": "internal_asset",
                        "chunk_index": idx,
                        "city": asset.get('city', 'Unknown'),
                        "asset_name": asset.get('name', 'Unknown')
                    })

                if asset_texts:
                    # Ingest into "deal" Namespace (Private Context)
                    print(f"Ingesting {len(asset_texts)} assets into 'deal' namespace...")
                    ingest_deal_assets(asset_texts, asset_metadatas, namespace="deal")
                
                # --- 2. Process Comps (Market) -> Namespace: market_comps ---
                comps = structured_data.get("comps", [])
                comp_texts = []
                comp_metadatas = []
                
                for idx, comp in enumerate(comps):
                    text_blob = f"""
                    Comparable Asset: {comp.get('name', 'Unknown')}
                    Type: {comp.get('asset_type', 'Logistics')}
                    Size: {comp.get('size_m2', 0)} sqm
                    Rent: {comp.get('rent_psm_pa', 0)} /sqm
                    Yield: {comp.get('yield', 'N/A')}
                    Date: {comp.get('acquisition_date', '')}
                    Notes: {comp.get('notes', '')}
                    """
                    comp_texts.append(text_blob.strip())
                    comp_metadatas.append({
                        "source": json_files[0],
                        "deal_id": deal_id, 
                        "record_type": "market_comp",
                        "chunk_index": idx,
                        "city": "Unknown" 
                    })
                
                if comp_texts:
                    # Ingest into "market_comps" Namespace (Public Knowledge Base)
                    print(f"Ingesting {len(comp_texts)} comps into 'market_comps' namespace...")
                    ingest_deal_assets(comp_texts, comp_metadatas, namespace="market_comps")

            except Exception as e:
                print(f"Error reading JSON or Ingesting to Vector DB: {e}")
                return {
                    "messages": [AIMessage(content=f"Error processing JSON/Vector DB: {e}", name="system_log")]
                }

    msg = f"Loaded structured data from {json_files[0]}" if json_files else "No structured JSON found."
    
    return {
        "messages": [AIMessage(content=msg, name="system_log")],
        "extracted_data": {"source_json": structured_data}
    }

def load_pdf_documents(state: DealState):
    """
    Step 3: Load PDF Documents
    Parses the PDF files.
    """
    print("--- Node: Load PDF Documents ---")
    
    base_dir = os.getcwd()
    if "backend" in base_dir:
        if base_dir.endswith("backend"):
             data_root = os.path.join(base_dir, "data")
        else:
             current_file_dir = os.path.dirname(os.path.abspath(__file__))
             backend_dir = os.path.dirname(os.path.dirname(current_file_dir))
             data_root = os.path.join(backend_dir, "data")
    else:
        data_root = os.path.join(base_dir, "backend", "data")
        
    pdf_dir = os.path.join(data_root, "raw_pdfs")
    pdf_texts = []
    pdf_files = []
    
    if os.path.exists(pdf_dir):
        pdf_files = [f for f in os.listdir(pdf_dir) if f.endswith('.pdf')]
        for pdf_file in pdf_files:
            try:
                full_path = os.path.join(pdf_dir, pdf_file)
                # Send a log for each file being parsed (optional, or just one summary log)
                # Since we can't stream partial updates easily within a node without custom callback, 
                # we'll just do the work and return the result.
                text = parse_pdf_document.invoke(full_path)
                pdf_texts.append(f"--- Document: {pdf_file} ---\n{text[:10000]}...\n")
            except Exception as e:
                print(f"Error reading PDF {pdf_file}: {e}")

    msg = f"Parsed {len(pdf_files)} PDF documents: {', '.join(pdf_files)}"
    
    # Merge with existing extracted_data
    existing_data = state.get("extracted_data", {})
    existing_data["pdf_texts"] = pdf_texts
    existing_data["pdf_files"] = pdf_files

    return {
        "messages": [AIMessage(content=msg, name="system_log")],
        "extracted_data": existing_data
    }

def align_with_llm(state: DealState):
    """
    Step 4: Align with LLM
    Uses LLM to align JSON and PDF data.
    """
    print("--- Node: Align with LLM ---")
    
    extracted = state.get("extracted_data", {})
    structured_data = extracted.get("source_json", {})
    pdf_texts = extracted.get("pdf_texts", [])
    pdf_files = extracted.get("pdf_files", [])

    if not structured_data and not pdf_texts:
        return {
            "messages": [
                AIMessage(content="No data to align.", name="system_log"),
                AIMessage(content="I could not find any data to process.", name="agent")
            ]
        }

    llm = ChatOpenAI(model="gpt-4o", temperature=0.2)

    prompt = f"""
    You are an expert Real Estate Investment Analyst.

    I have two data sources:
    1. **Structured Data (JSON)**: This is the source of truth for IDs and basic specs.
    {json.dumps(structured_data, indent=2)[:5000]} 

    2. **Unstructured Data (PDF Content)**:
    {"".join(pdf_texts)}

    **Your Task:**
    1. **Align**: Match the assets in the JSON with the descriptions in the PDF. Use Address or Property Name as the key.
    2. **Extract & Enrich**: For each matched asset, extract the following from the PDF and add it to the data:
       - 'market_highlights': Key selling points of the location.
       - 'investment_rationale': Why is this a good deal?
       - 'risk_factors': Any mentioned risks.
       - 'financials': Extract NOI, ERV, Cap Rate if available.
       - 'physical_specs': Extract Clear Height, Floor Loading, Dock Doors if available.
    3. **Verify**: Check if the GLA (Gross Leasable Area) and Occupancy in the JSON match the PDF. If different, create a 'discrepancies' field.

    Return a summary of the aligned deal data in a readable format for the user, highlighting the key metrics and any discrepancies found.
    """

    response = llm.invoke([HumanMessage(content=prompt)])

    status_msg = "I’ve ingested the IM, rent roll and structured deal data.\n\nI’ll generate the summary, key metrics, and an initial set of comparables."
    
    # Update extracted data with analysis
    extracted["analysis"] = response.content

    return {
        "messages": [
            AIMessage(content=status_msg, name="system_log"),
            AIMessage(content=response.content, name="agent")
        ],
        "extracted_data": extracted
    }

def compute_metrics_and_draft_summary(state: DealState):
    """
    Step 5: Compute Metrics and Draft Summary
    Reads initial data and computes preliminary summaries.
    """
    print("--- Node: Compute Metrics and Draft Summary ---")

    extracted = state.get("extracted_data", {})
    source_json = extracted.get("source_json", {})
    analysis_text = extracted.get("analysis", "")

    if not source_json and not analysis_text:
        return {"messages": [AIMessage(content="No data available to compute metrics.", name="agent")]}

    # Use LLM to synthesize metrics from the ingested data
    llm = ChatOpenAI(model="gpt-4o", temperature=0.2)

    prompt = f"""
    You are a Real Estate Analyst. 
    Based on the provided structured data and the analysis from the ingestion step, generate a concise "Metrics Summary".

    **Source Data (JSON):**
    {json.dumps(source_json, indent=2)[:2000]}

    **Preliminary Analysis:**
    {analysis_text}

    **Task:**
    Compute or extract the following key metrics. Keep it brief (bullet points).

    **Output Format:**
    Compute Metrics and Draft Summary:
    - Total GLA: [Value]
    - Occupancy: [Value]
    - WALT: [Value]
    - In-Place Rent: [Value]
    - Key Highlight 1: [Brief text]
    - Key Highlight 2: [Brief text]
    """

    response = llm.invoke([HumanMessage(content=prompt)])
    content = response.content

    # --- Generate PPT ---
    ppt_link_msg = ""
    try:
        # Paths
        current_dir = os.path.dirname(os.path.abspath(__file__))
        backend_dir = os.path.dirname(os.path.dirname(current_dir))
        # Use the specific Deal Summary template
        template_path = os.path.join(backend_dir, "data", "templates", "deal_summary_template.pptx")
        output_dir = os.path.join(backend_dir, "data", "generated")
        os.makedirs(output_dir, exist_ok=True)

        # Load Template
        if os.path.exists(template_path):
            try:
                prs = Presentation(template_path)
            except:
                prs = Presentation()
        else:
            prs = Presentation()
            # Fallback layout
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "{{DEAL_NAME}}"
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "{{DATE}}"
            
            if len(prs.slide_layouts) > 1:
                slide_layout = prs.slide_layouts[1]
            else:
                slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Executive Summary"
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "{{SUMMARY_BULLETS}}"

        # Replacements
        replacements = {
            "{{DEAL_NAME}}": state.get("company_name", "Project Deal") or "Project Deal",
            "{{DATE}}": datetime.now().strftime("%Y-%m-%d"),
            "{{SUMMARY_BULLETS}}": content,
            # Deal Summary might not have these yet, but we keep them for safety
            "{{MARKET_BULLETS}}": "Pending Market Analysis...",
            "{{ENTRY_YIELD}}": "TBD",
            "{{IRR}}": "TBD",
            "{{MOIC}}": "TBD"
        }

        # Helper to replace text
        def replace_text(text_frame, replacements):
            for paragraph in text_frame.paragraphs:
                full_text = paragraph.text
                original_text = full_text
                for key, val in replacements.items():
                    if key in full_text:
                        full_text = full_text.replace(key, str(val))
                if full_text != original_text:
                    paragraph.text = full_text

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    replace_text(shape.text_frame, replacements)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                replace_text(cell.text_frame, replacements)

        # Save locally
        ppt_filename = f"Deal_Summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = os.path.join(output_dir, ppt_filename)
        
        prs.save(output_path)
        
        # Upload to S3
        s3_link = upload_to_s3_and_get_link(output_path)
        
        # Remove local
        if s3_link and os.path.exists(output_path):
            os.remove(output_path)
        
        if s3_link:
            ppt_link_msg = f"\n\n\n📥 **[Download Deal Summary (PPT)]({s3_link})**"
        else:
            ppt_link_msg = f"\n\n(PPT generated locally at {output_path}, but S3 upload failed - check AWS credentials)"
            
    except Exception as e:
        print(f"PPT Generation Error: {e}")
        ppt_link_msg = f"\n\n(Error generating PPT: {e})"

    return {"messages": [AIMessage(content=content + ppt_link_msg, name="agent")]}
